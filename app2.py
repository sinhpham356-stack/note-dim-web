import streamlit as st
from PIL import Image, ImageDraw, ImageFont
import io
import base64
from pptx import Presentation
from pptx.util import Inches

# --- BẢN VÁ LỖI MÀN HÌNH ĐEN CHO STREAMLIT CLOUD ---
import streamlit.elements.image as st_image
def patched_image_to_url(image, *args, **kwargs):
    buffered = io.BytesIO()
    if isinstance(image, Image.Image):
        # Bắt buộc chuyển sang RGB để không bị đen nền
        image.convert("RGB").save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()
    return f"data:image/png;base64,{img_str}"
st_image.image_to_url = patched_image_to_url

from streamlit_drawable_canvas import st_canvas
# ------------------------------------------------

st.set_page_config(page_title="Note-Dim Web", layout="wide")

if 'project_slides' not in st.session_state:
    st.session_state.project_slides = []

def create_pdf(images):
    buf = io.BytesIO()
    if images:
        rgb_images = [img.convert('RGB') for img in images]
        rgb_images[0].save(buf, format='PDF', save_all=True, append_images=rgb_images[1:])
    buf.seek(0)
    return buf.getvalue()

def create_pptx(images):
    prs = Presentation()
    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img_io = io.BytesIO()
        img.save(img_io, format='PNG')
        img_io.seek(0)
        slide.shapes.add_picture(img_io, Inches(0), Inches(0), width=prs.slide_width)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

# --- SIDEBAR: QUẢN LÝ SLIDE BÁO CÁO ---
with st.sidebar:
    st.header("📑 Danh sách Trang (Slide)")
    st.write(f"Đang có: **{len(st.session_state.project_slides)}** ảnh")
    
    if len(st.session_state.project_slides) > 0:
        for i, img in enumerate(st.session_state.project_slides):
            st.image(img, caption=f"Trang {i+1}", use_container_width=True)
            
        st.write("---")
        st.download_button("📄 Tải xuống PDF", data=create_pdf(st.session_state.project_slides), file_name="Bao_Cao_Ky_Thuat.pdf", mime="application/pdf")
        st.download_button("📊 Tải xuống PowerPoint", data=create_pptx(st.session_state.project_slides), file_name="Bao_Cao_Ky_Thuat.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        
        if st.button("🗑️ Xóa toàn bộ dự án"):
            st.session_state.project_slides = []
            st.rerun()

# --- MAIN APP ---
st.title("📏 Công cụ Note-Dim (Phiên bản Báo cáo Web)")

uploaded_file = st.file_uploader("📂 Tải ảnh bản vẽ lên", type=["png", "jpg", "jpeg"])

if uploaded_file is not None:
    # Xử lý ảnh gốc để đảm bảo không bị lỗi màu
    original_image = Image.open(uploaded_file).convert("RGB")
    
    st.write("---")
    st.markdown("### Bước 1: Vẽ các đường nét (Thay thế Dim/Góc)")
    col1, col2, col3 = st.columns([2, 2, 2])
    with col1:
        drawing_mode = st.selectbox("🖱️ Chọn công cụ:", ("line", "freedraw", "rect", "polygon"))
    with col2:
        stroke_width = st.slider("📐 Độ dày nét:", 1, 10, 3)
    with col3:
        stroke_color = st.color_picker("🎨 Màu nét vẽ:", "#FF0000")

    canvas_width = 1000
    canvas_height = int((original_image.height / original_image.width) * canvas_width)
    
    # BẢNG VÊ CAN-VAS
    canvas_result = st_canvas(
        fill_color="rgba(0, 0, 0, 0)",
        stroke_width=stroke_width,
        stroke_color=stroke_color,
        background_image=original_image,
        update_streamlit=True,
        width=canvas_width,
        height=canvas_height,
        drawing_mode=drawing_mode,
        key="canvas",
    )

    if canvas_result.image_data is not None:
        # Gộp nét vẽ vào ảnh gốc
        drawn_layer = Image.fromarray(canvas_result.image_data.astype("uint8"), "RGBA")
        drawn_layer = drawn_layer.resize(original_image.size, Image.Resampling.LANCZOS)
        base_img_with_drawing = original_image.convert("RGBA").copy()
        base_img_with_drawing.alpha_composite(drawn_layer)
        final_img = base_img_with_drawing.convert("RGB")

        st.write("---")
        st.markdown("### Bước 2: Thêm Tag/Note Text (Do web không gõ trực tiếp được)")
        
        c1, c2, c3 = st.columns([3, 1, 1])
        with c1:
            note_text = st.text_input("📝 Gõ nội dung (Kích thước, mã thiết bị, v.v.):")
        with c2:
            text_size = st.number_input("Cỡ chữ:", min_value=10, max_value=200, value=40)
        with c3:
            text_color = st.color_picker("Màu chữ:", "#FFFF00") # Mặc định Vàng

        if note_text:
            st.info("Kéo thanh trượt bên dưới để di chuyển vị trí đoạn Text trên ảnh:")
            pos_x = st.slider("↔️ Căn ngang (X):", 0, original_image.width, original_image.width // 2)
            pos_y = st.slider("↕️ Căn dọc (Y):", 0, original_image.height, original_image.height // 2)
            
            # Đóng dấu Text lên ảnh
            draw = ImageDraw.Draw(final_img)
            try:
                # Dùng font mặc định của Pillow
                font = ImageFont.load_default()
                # Hack tăng size cho font mặc định
                final_img_temp = final_img.copy()
                draw_temp = ImageDraw.Draw(final_img_temp)
                # Vẽ bóng mờ đen (outline) để chữ dễ đọc trên mọi nền
                draw_temp.text((pos_x+2, pos_y+2), note_text, fill="black")
                draw_temp.text((pos_x, pos_y), note_text, fill=text_color)
                final_img = final_img_temp
            except Exception:
                pass

        st.write("### 🖼️ XEM TRƯỚC BẢN IN")
        st.image(final_img, use_container_width=True)

        if st.button("➕ CHỐT: Thêm ảnh này vào Báo cáo (Slide)", type="primary"):
            st.session_state.project_slides.append(final_img)
            st.success(f"Đã thêm! Tổng số lượng: {len(st.session_state.project_slides)} trang.")
            st.rerun()
